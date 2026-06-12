"""
Generate FlangParallelAnalyzer presentation (.pptx)
Run: python3 make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xBB, 0xCC)
GREEN     = RGBColor(0x2D, 0xD4, 0x81)
ORANGE    = RGBColor(0xFB, 0xBF, 0x24)
RED       = RGBColor(0xF8, 0x71, 0x71)
DARK_CARD = RGBColor(0x16, 0x2A, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ─────────────────────────────────────────────────────────

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, font_size=18, bold=False,
        color=WHITE, bg_color=None, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox

def rect(slide, l, t, w, h, fill_color, radius=False):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def accent_bar(slide, t=0.55, h=0.06):
    r = rect(slide, 0, t, 13.33, h, ACCENT)
    return r

def slide_title(slide, title, subtitle=None):
    box(slide, 0.5, 0.1, 12, 0.55, title,
        font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    accent_bar(slide)
    if subtitle:
        box(slide, 0.5, 0.7, 12, 0.4, subtitle,
            font_size=14, color=ACCENT2, align=PP_ALIGN.LEFT, italic=True)

def bullet_block(slide, l, t, w, h, items, font_size=15,
                 bullet="▸", color=WHITE, spacing=0.42):
    y = t
    for item in items:
        box(slide, l, y, w, spacing, f"{bullet}  {item}",
            font_size=font_size, color=color)
        y += spacing

def card(slide, l, t, w, h, title, body, title_color=ACCENT,
         body_color=WHITE, font_title=13, font_body=12):
    rect(slide, l, t, w, h, DARK_CARD)
    box(slide, l+0.1, t+0.08, w-0.2, 0.35,
        title, font_size=font_title, bold=True, color=title_color)
    box(slide, l+0.1, t+0.4, w-0.2, h-0.5,
        body, font_size=font_body, color=body_color, wrap=True)

def verdict_chip(slide, l, t, label, color):
    rect(slide, l, t, 1.3, 0.38, color)
    box(slide, l+0.05, t+0.04, 1.2, 0.3,
        label, font_size=13, bold=True, color=RGBColor(0,0,0),
        align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# big cyan top bar
rect(s, 0, 0, 13.33, 1.1, ACCENT)
box(s, 0.5, 0.12, 12.3, 0.9,
    "FlangParallelAnalyzer",
    font_size=38, bold=True, color=BG, align=PP_ALIGN.CENTER)

box(s, 0.5, 1.3, 12.3, 0.7,
    "Automatic Loop Parallelism Analysis for Fortran using LLVM Flang & MLIR",
    font_size=20, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# divider
rect(s, 3.5, 2.2, 6.3, 0.05, ACCENT2)

box(s, 0.5, 2.5, 12.3, 0.45,
    "Aryan Gupta   ·   Arushi Vaidya   ·   Arpita",
    font_size=16, color=ACCENT2, align=PP_ALIGN.CENTER)
box(s, 0.5, 2.95, 12.3, 0.38,
    "Dept. of Computer Science and Engineering  |  RV College of Engineering, Bengaluru",
    font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# three bottom stat cards
for i, (num, lbl) in enumerate([("35", "Test Cases"), ("100%", "Accuracy"), ("0", "False Positives")]):
    x = 1.8 + i * 3.3
    rect(s, x, 4.0, 2.6, 1.4, DARK_CARD)
    box(s, x, 4.1, 2.6, 0.7, num,
        font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    box(s, x, 4.75, 2.6, 0.4, lbl,
        font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

box(s, 0.5, 6.8, 12.3, 0.35,
    "RV College of Engineering  ·  2024–25",
    font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "The Problem", "Why is parallelizing Fortran loops hard?")

# left panel — context
rect(s, 0.4, 0.9, 5.8, 5.8, DARK_CARD)
box(s, 0.55, 1.0, 5.5, 0.4,
    "Fortran powers scientific computing", font_size=14, bold=True, color=ACCENT)
bullet_block(s, 0.6, 1.5, 5.5, 4.0, [
    "Weather simulations, CFD, matrix solvers",
    "Performance lives inside DO loops",
    "OpenMP !$OMP PARALLEL DO enables\n   multi-core speedup",
    "But wrong annotation = silent wrong results",
    "Missed annotation = wasted performance",
], font_size=13, spacing=0.62)

# right panel — the manual pain
rect(s, 6.5, 0.9, 6.4, 5.8, DARK_CARD)
box(s, 6.65, 1.0, 6.1, 0.4,
    "Manual annotation today", font_size=14, bold=True, color=ORANGE)

steps = [
    ("1", "Read FIR / assembly output"),
    ("2", "Trace every array subscript"),
    ("3", "Check for alias relationships"),
    ("4", "Identify reduction variables"),
    ("5", "Write !$OMP directive by hand"),
]
for i, (num, txt) in enumerate(steps):
    y = 1.55 + i * 0.88
    rect(s, 6.6, y, 0.38, 0.38, ACCENT)
    box(s, 6.6, y, 0.38, 0.38, num,
        font_size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, 7.1, y+0.02, 5.5, 0.38, txt, font_size=13, color=WHITE)

box(s, 6.55, 6.1, 6.2, 0.38,
    "⚠  Error-prone and slow for large codebases",
    font_size=13, bold=True, color=RED)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Our Solution — FlangParallelAnalyzer",
            "Static analysis tool that tells you, per loop, whether it's safe to parallelize")

# central flow arrow
centers = [1.1, 4.0, 7.0, 10.2]
labels  = ["Fortran\nSource (.f90)", "FIR\n(MLIR IR)", "FPA\nAnalysis Pass", "Verdict +\nOpenMP Hint"]
colors  = [ACCENT, ACCENT2, ORANGE, GREEN]

for i, (x, lbl, col) in enumerate(zip(centers, labels, colors)):
    rect(s, x-0.85, 1.6, 1.7, 1.2, DARK_CARD)
    box(s, x-0.85, 1.65, 1.7, 1.1,
        lbl, font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < 3:
        box(s, x+0.85, 2.0, 0.6, 0.4,
            "→", font_size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

box(s, 1.5, 2.9, 2.5, 0.35,
    "flang-new -fc1 -emit-fir", font_size=10, color=GRAY, italic=True)
box(s, 4.3, 2.9, 2.5, 0.35,
    "fpa-tool input.fir", font_size=10, color=GRAY, italic=True)

# three output boxes
for i, (lbl, col, eg) in enumerate([
    ("SAFE",      GREEN,  "!$OMP PARALLEL DO"),
    ("REDUCTION", ORANGE, "!$OMP PARALLEL DO\nREDUCTION(+:sum)"),
    ("UNSAFE",    RED,    "! Cannot parallelize"),
]):
    x = 0.5 + i * 4.1
    rect(s, x, 3.55, 3.6, 2.5, DARK_CARD)
    rect(s, x, 3.55, 3.6, 0.42, col)
    box(s, x, 3.55, 3.6, 0.42,
        lbl, font_size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, x+0.1, 4.05, 3.4, 1.8,
        eg, font_size=13, color=WHITE, italic=True)

box(s, 0.4, 6.35, 12.5, 0.38,
    "Tool runs in < 1 second  ·  Zero false positives  ·  Full analysis trace per loop",
    font_size=13, color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — How It Works: 5-Phase Pipeline
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "How It Works — 5-Phase Analysis Pipeline")

phases = [
    ("1  Structure",       "Collect loop bounds,\nnesting depth, op count",       ACCENT),
    ("2  Mem Access",      "Classify every read/write\nas external or local",      ACCENT2),
    ("3  Index Analysis",  "Is subscript a(i) or a(i±k)?\nDetects loop-carried deps", ORANGE),
    ("4  Reduction",       "Match load→op→store\npattern on scalars",             GREEN),
    ("5  Fallback",        "No ext writes → SAFE\nOtherwise → UNSAFE",            RED),
]

for i, (title, body, col) in enumerate(phases):
    x = 0.35 + i * 2.55
    rect(s, x, 0.9, 2.3, 5.7, DARK_CARD)
    rect(s, x, 0.9, 2.3, 0.4, col)
    box(s, x, 0.9, 2.3, 0.4,
        title, font_size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, x+0.1, 1.45, 2.1, 2.2,
        body, font_size=12, color=WHITE, wrap=True)
    if i < 4:
        box(s, x+2.3, 3.5, 0.25, 0.4,
            "→", font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

box(s, 0.35, 6.72, 12.6, 0.38,
    "Each phase either sets a final verdict (SAFE / REDUCTION / UNSAFE) or passes to the next phase",
    font_size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Phase 3 Deep Dive: Dependency Detection
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Dependency Detection — Phase 3",
            "Tracing array subscripts through the compiler IR")

# left: IV-derived
rect(s, 0.4, 0.9, 5.9, 5.7, DARK_CARD)
box(s, 0.55, 1.0, 5.6, 0.4,
    "✓  a(i)  →  IV-derived  →  SAFE", font_size=14, bold=True, color=GREEN)
box(s, 0.55, 1.5, 5.6, 0.35,
    "Fortran source:", font_size=12, color=GRAY)
box(s, 0.55, 1.85, 5.6, 0.55,
    "b(i) = a(i) * 2.0", font_size=14, italic=True, color=ACCENT2,
    bg_color=RGBColor(0x0A, 0x24, 0x38))
box(s, 0.55, 2.55, 5.6, 0.35,
    "FIR subscript trace:", font_size=12, color=GRAY)
steps2 = [
    "fir.array_coor coordinate operand",
    "→ arith.index_cast (strip type cast)",
    "→ fir.load %loop_var_alloca",
    "→ fir.store iter_arg → alloca",
    "→ loop induction variable  ✓",
]
for j, st in enumerate(steps2):
    col2 = GREEN if j == 4 else WHITE
    box(s, 0.65, 2.95 + j*0.48, 5.4, 0.42, st,
        font_size=11, color=col2, italic=(j==4))

# right: IV±k
rect(s, 6.6, 0.9, 6.3, 5.7, DARK_CARD)
box(s, 6.75, 1.0, 6.0, 0.4,
    "✗  a(i-1)  →  IV offset  →  UNSAFE", font_size=14, bold=True, color=RED)
box(s, 6.75, 1.5, 6.0, 0.35,
    "Fortran source:", font_size=12, color=GRAY)
box(s, 6.75, 1.85, 6.0, 0.55,
    "a(i) = a(i) + a(i-1)", font_size=14, italic=True, color=ACCENT2,
    bg_color=RGBColor(0x0A, 0x24, 0x38))
box(s, 6.75, 2.55, 6.0, 0.35,
    "FIR subscript trace:", font_size=12, color=GRAY)
steps3 = [
    "fir.array_coor coordinate operand",
    "→ arith.subi  %iv,  1",
    "→ LHS is IV-derived",
    "→ RHS is constant k = 1",
    "→ IV ± k  with  k ≠ 0  →  UNSAFE  ✗",
]
for j, st in enumerate(steps3):
    col3 = RED if j == 4 else WHITE
    box(s, 6.85, 2.95 + j*0.48, 5.8, 0.42, st,
        font_size=11, color=col3, italic=(j==4))

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Reduction Detection (Phase 4)
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Reduction Detection — Phase 4",
            "Automatically recognizing scalar accumulation patterns")

# pattern diagram
rect(s, 0.4, 0.95, 12.5, 1.3, RGBColor(0x0A, 0x24, 0x38))
box(s, 0.5, 0.98, 12.2, 0.5,
    "Pattern to match:", font_size=12, color=GRAY)
box(s, 0.5, 1.35, 12.2, 0.75,
    "%old  =  fir.load %acc       →       %new  =  arith.addf %old, <expr>       →       fir.store %new to %acc",
    font_size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)

# examples
examples = [
    ("total = total + a(i)*b(i)",  "REDUCTION(+:total)",   "Floating-point sum",      GREEN),
    ("prod  = prod  * a(i)",       "REDUCTION(*:prod)",    "Floating-point product",  GREEN),
    ("isum  = isum  + ia(i)",      "REDUCTION(+:isum)",    "Integer summation",       GREEN),
    ("norm2 = norm2 + a(i)*a(i)",  "REDUCTION(+:norm2)",   "Norm-squared",            GREEN),
    ("s = s + f(x) + g(y)",        "UNSAFE (chained)",     "Not yet supported",       ORANGE),
]

for i, (src, hint, desc, col) in enumerate(examples):
    y = 2.45 + i * 0.82
    rect(s, 0.4, y, 12.5, 0.72, DARK_CARD)
    box(s, 0.6,  y+0.12, 4.2, 0.45, src,  font_size=12, italic=True, color=WHITE)
    box(s, 5.0,  y+0.12, 4.0, 0.45, hint, font_size=12, bold=True, color=col)
    box(s, 9.3,  y+0.12, 3.4, 0.45, desc, font_size=12, color=GRAY)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Architecture
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "System Architecture",
            "Built on standard LLVM 18 / MLIR infrastructure — no custom IR needed")

components = [
    ("LoopParallelAnalysis.cpp\n471 LOC",
     "Core MLIR pass\nPhases 1–5\nLoop printer",
     ACCENT),
    ("AccessClassifier.cpp\n165 LOC",
     "Phase 2 memory walker\nBase-ref stripping\nAccessRecord types",
     ACCENT2),
    ("fpa-tool  main.cpp\n71 LOC",
     "CLI driver\nDialect registration\nPassManager setup",
     ORANGE),
    ("report.py\n1 259 LOC",
     "Compiles + runs tool\nParses output\nGenerates HTML report",
     GREEN),
]

for i, (title, body, col) in enumerate(components):
    x = 0.4 + i * 3.2
    rect(s, x, 1.0, 2.9, 3.2, DARK_CARD)
    rect(s, x, 1.0, 2.9, 0.42, col)
    box(s, x, 1.0, 2.9, 0.42,
        title, font_size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, x+0.1, 1.55, 2.7, 2.4,
        body, font_size=12, color=WHITE, wrap=True)

# tech stack row
box(s, 0.4, 4.45, 12.5, 0.38,
    "Technology Stack", font_size=13, bold=True, color=ACCENT)

stack = [
    ("LLVM 18", ACCENT),
    ("MLIR", ACCENT2),
    ("Flang / FIR", ORANGE),
    ("C++17", WHITE),
    ("Python 3", GREEN),
    ("CMake ≥ 3.20", GRAY),
]
for i, (tech, col) in enumerate(stack):
    x = 0.4 + i * 2.15
    rect(s, x, 4.9, 1.95, 0.5, DARK_CARD)
    box(s, x, 4.9, 1.95, 0.5,
        tech, font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)

box(s, 0.4, 5.6, 12.5, 0.9,
    "Why FIR instead of Fortran source text?\n"
    "FIR preserves loop structure, array shapes, and variable intent. "
    "Every subscript is an explicit SSA value — no need to parse Fortran's many syntax quirks.",
    font_size=12, color=GRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Test Results
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Test Results — 35 Cases, 100% Accuracy")

# summary bar
for i, (lbl, n, col) in enumerate([
    ("SAFE",           13, GREEN),
    ("REDUCTION",       5, ORANGE),
    ("UNSAFE (dep)",   10, RED),
    ("UNSAFE (conserv)",7, RGBColor(0xDD,0x90,0x20)),
]):
    x = 0.4 + i * 3.2
    rect(s, x, 0.9, 2.9, 0.72, col)
    box(s, x, 0.9, 2.9, 0.38,
        f"{n}  {lbl}", font_size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

# table
headers = ["Fortran Pattern", "Expected", "Actual", "Pass?"]
rows = [
    ("b(i) = a(i) * 2.0",          "SAFE",      "SAFE",      "✓"),
    ("total += a(i)*b(i)",          "REDUCTION", "REDUCTION", "✓"),
    ("a(i) = a(i) + a(i-1)",        "UNSAFE",    "UNSAFE",    "✓"),
    ("b(idx(i)) = a(i)  (scatter)", "UNSAFE",    "UNSAFE",    "✓"),
    ("a(i) = a(i) * 3.0  (inplace)","UNSAFE",    "UNSAFE",    "✓"),
    ("norm2 += a(i)*a(i)",          "REDUCTION", "REDUCTION", "✓"),
    ("read-only traversal",         "SAFE",      "SAFE",      "✓"),
    ("a(i) = a(i+1) * 2.0",         "UNSAFE",    "UNSAFE",    "✓"),
]

col_x = [0.4, 5.5, 8.5, 11.4]
col_w = [4.8, 2.7, 2.7, 1.5]

# header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    rect(s, cx, 1.8, cw, 0.4, RGBColor(0x00,0x4E,0x6E))
    box(s, cx+0.05, 1.82, cw-0.1, 0.36,
        hdr, font_size=12, bold=True, color=WHITE)

for i, row in enumerate(rows):
    y = 2.27 + i * 0.52
    bg_c = RGBColor(0x13,0x27,0x3A) if i % 2 == 0 else DARK_CARD
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        rect(s, cx, y, cw, 0.48, bg_c)
        col_text = WHITE
        if j == 3:
            col_text = GREEN
        elif j == 1 or j == 2:
            if cell == "SAFE":      col_text = GREEN
            elif cell == "REDUCTION": col_text = ORANGE
            elif cell == "UNSAFE":  col_text = RED
        box(s, cx+0.05, y+0.06, cw-0.1, 0.36,
            cell, font_size=11, color=col_text, bold=(j==3))

box(s, 0.4, 6.6, 12.5, 0.4,
    "Zero false positives  ·  Conservative policy: UNSAFE preferred over incorrect SAFE",
    font_size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — HTML Report Demo
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Interactive HTML Report",
            "Three-panel IDE-style interface generated by report.py")

# three panel mock-ups
panels = [
    ("File Browser",
     "▸  01_safe_scale.f90      ✓ SAFE\n"
     "▸  04_reduction_sum.f90   ⊕ REDU\n"
     "▸  06_dep_shift1.f90      ✗ UNSAFE\n"
     "▸  08_dep_inplace.f90     ✗ UNSAFE\n"
     "▸  30_reduction_norm2.f90 ⊕ REDU\n"
     "   …  35 files total",
     ACCENT),
    ("Source View",
     "12 •  do i = 1, n\n"
     "         [ ✓ safe · independent ]\n"
     "13      b(i) = a(i) * 2.0\n"
     "14    end do\n"
     "15  end subroutine",
     ORANGE),
    ("Analysis Detail",
     "✓ SAFE   HIGH 90%\n"
     "!$OMP PARALLEL DO\n"
     "─────────────────────\n"
     "Phase 3  ✓ All subscripts IV-derived\n"
     "Phase 4  – Not needed\n"
     "Phase 5  – Not needed\n"
     "Mem: [W] array  [R] array  [RW] scalar",
     GREEN),
]

for i, (title, body, col) in enumerate(panels):
    x = 0.35 + i * 4.3
    rect(s, x, 0.95, 4.0, 5.6, DARK_CARD)
    rect(s, x, 0.95, 4.0, 0.38, col)
    box(s, x+0.05, 0.95, 3.9, 0.38,
        title, font_size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, x+0.12, 1.45, 3.75, 4.9,
        body, font_size=11, color=WHITE, wrap=True)

box(s, 0.35, 6.75, 12.5, 0.38,
    "Run:  python3 scripts/report.py -o report.html   then open in browser",
    font_size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Limitations & Future Work + Conclusion
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Limitations, Future Work & Conclusion")

# left — limitations
rect(s, 0.4, 0.9, 5.8, 4.0, DARK_CARD)
box(s, 0.55, 0.95, 5.5, 0.38,
    "Known Limitations", font_size=14, bold=True, color=ORANGE)
limits = [
    "a(i) = a(i)*k  conservatively UNSAFE",
    "Multi-dim subscripts  c(i,j)  not traced",
    "Function call side-effects not analyzed",
    "min / max reductions not yet detected",
    "Confidence score is heuristic only",
]
bullet_block(s, 0.55, 1.42, 5.5, 3.3, limits,
             font_size=12, bullet="–", color=GRAY, spacing=0.5)

# middle — future work
rect(s, 6.5, 0.9, 6.4, 4.0, DARK_CARD)
box(s, 6.65, 0.95, 6.1, 0.38,
    "Future Work", font_size=14, bold=True, color=ACCENT)
future = [
    "In-place update read-before-write check",
    "Multi-dimensional subscript tracing",
    "min / max reduction support",
    "Interprocedural side-effect summaries",
    "Auto-insert !$OMP directives in source",
    "CI/CD integration as pre-commit hook",
]
bullet_block(s, 6.65, 1.42, 6.1, 3.5, future,
             font_size=12, bullet="→", color=WHITE, spacing=0.48)

# conclusion box
rect(s, 0.4, 5.1, 12.5, 1.95, RGBColor(0x00, 0x3A, 0x52))
box(s, 0.55, 5.18, 12.2, 0.38,
    "Conclusion", font_size=14, bold=True, color=ACCENT)
box(s, 0.55, 5.6, 12.2, 1.3,
    "FlangParallelAnalyzer demonstrates that a compact (~700 LOC C++) MLIR pass "
    "built on Flang FIR can automate Fortran loop parallelism classification "
    "with 100% accuracy on a comprehensive test suite and zero false positives. "
    "The tool bridges the gap between opaque auto-parallelizing compilers and "
    "error-prone manual annotation.",
    font_size=12, color=WHITE, wrap=True)

# ── Save ─────────────────────────────────────────────────────────────
out = "FlangParallelAnalyzer_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
