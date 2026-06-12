"""
Generate FlangParallelAnalyzer_v2.pptx — matched to reference CD_Lab_EL.pptx format.

Run: python3 paper/make_ppt_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os, zipfile, shutil, re

# ── Load reference ───────────────────────────────────────────────────────────
REF_PATH = "/tmp/ref.pptx"
OUT_PATH = "FlangParallelAnalyzer_v2.pptx"

ref = Presentation(REF_PATH)

# Extract image bytes so we can embed them in a fresh presentation
with zipfile.ZipFile(REF_PATH) as z:
    img1_bytes = z.read("ppt/media/image1.png")   # background texture
    try:
        img2_bytes = z.read("ppt/media/image2.jpeg")
    except Exception:
        img2_bytes = None

# ── Create fresh presentation with same dimensions ───────────────────────────
prs = Presentation()
prs.slide_width  = ref.slide_width    # 20 inches
prs.slide_height = ref.slide_height   # 11.25 inches

BLANK = prs.slide_layouts[6]

# ── Colours / constants ───────────────────────────────────────────────────────
BLUE   = "083C92"
BLACK  = "000000"
WHITE  = "FFFFFF"
GRAY   = "888888"
LGRAY  = "D9D9D9"

TITLE_FONT  = "Playfair Display Bold Italics"
BODY_FONT   = "Calibri"

W  = 20.0   # slide width  (inches)
H  = 11.25  # slide height (inches)

# ── XML namespace shortcuts ───────────────────────────────────────────────────
nsmap = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

def _emu(inches): return int(inches * 914400)

# ── Build the background blipFill freeform XML string ────────────────────────
def _bg_freeform_xml(rId):
    w = _emu(W); h = _emu(H)
    return f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr name="Background" id="2"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm flipH="false" flipV="false" rot="0">
      <a:off x="0" y="0"/>
      <a:ext cx="{w}" cy="{h}"/>
    </a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect r="r" b="b" t="t" l="l"/>
      <a:pathLst>
        <a:path h="{h}" w="{w}">
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="{w}" y="0"/></a:lnTo>
          <a:lnTo><a:pt x="{w}" y="{h}"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="{h}"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="0"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:blipFill>
      <a:blip r:embed="{rId}"/>
      <a:stretch><a:fillRect l="0" t="0" r="0" b="0"/></a:stretch>
    </a:blipFill>
  </p:spPr>
</p:sp>"""

# ── Add background image to a slide and return its relationship id ────────────
def add_bg(slide):
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    img_part, rId = slide.part.get_or_add_image_part(
        __import__("io").BytesIO(img1_bytes)
    )
    return rId

# ── Low-level helper: add XML element to slide spTree ─────────────────────────
def add_xml(slide, xml_str):
    el = etree.fromstring(xml_str)
    slide.shapes._spTree.append(el)
    return el

# ── Textbox XML builder ───────────────────────────────────────────────────────
def _txb_xml(shape_id, x, y, w, h, paragraphs, no_inset=True):
    """
    paragraphs: list of dicts with keys:
      text, sz (half-pts), bold, italic, color, align, font, space_before, line_space
    """
    inset = 'tIns="0" lIns="0" bIns="0" rIns="0"' if no_inset else ""
    p_xml_parts = []
    for i, pg in enumerate(paragraphs):
        align_map = {PP_ALIGN.CENTER: "ctr", PP_ALIGN.LEFT: "l",
                     PP_ALIGN.RIGHT: "r", None: "l"}
        algn    = align_map.get(pg.get("align"), "l")
        sb      = pg.get("space_before", 0)
        ls      = pg.get("line_space", 0)
        spc_xml = ""
        if sb:
            spc_xml += f'<a:spcBef><a:spcPts val="{sb}"/></a:spcBef>'
        if ls:
            spc_xml += f'<a:lnSpc><a:spcPts val="{ls}"/></a:lnSpc>'

        text  = pg.get("text", "")
        sz    = pg.get("sz", 2000)
        bold  = "true" if pg.get("bold") else "false"
        ital  = "true" if pg.get("italic") else "false"
        color = pg.get("color", BLACK)
        font  = pg.get("font", BODY_FONT)

        # Split on \n for line breaks within a paragraph
        lines = text.split("\n")
        runs = []
        for j, line in enumerate(lines):
            if j > 0:
                runs.append("<a:br/>")
            if line:
                runs.append(f"""<a:r>
  <a:rPr lang="en-US" b="{bold}" i="{ital}" sz="{sz}" dirty="0">
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:latin typeface="{font}"/>
  </a:rPr>
  <a:t>{_esc(line)}</a:t>
</a:r>""")

        p_xml_parts.append(f"""<a:p>
  <a:pPr algn="{algn}">{spc_xml}</a:pPr>
  {"".join(runs)}
</a:p>""")

    return f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr name="TextBox {shape_id}" id="{shape_id}"/>
    <p:cNvSpPr txBox="true"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm rot="0">
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="t" rtlCol="false" {inset}>
      <a:spAutoFit/>
    </a:bodyPr>
    <a:lstStyle/>
    {"".join(p_xml_parts)}
  </p:txBody>
</p:sp>"""

def _esc(s):
    return (s.replace("&","&amp;").replace("<","&lt;")
             .replace(">","&gt;").replace('"',"&quot;"))

# ── White content-area rectangle ──────────────────────────────────────────────
def _content_rect_xml(shape_id, x, y, w, h, fill=WHITE, radius_emu=0):
    return f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr name="ContentBox {shape_id}" id="{shape_id}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm rot="0">
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="roundRect">
      <a:avLst><a:gd name="adj" fmla="val 15000"/></a:avLst>
    </a:prstGeom>
    <a:solidFill><a:srgbClr val="{fill}"><a:alpha val="92000"/></a:srgbClr></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
</p:sp>"""

# ── Blue banner for slide titles ───────────────────────────────────────────────
def _title_banner_xml(shape_id, x, y, w, h, text, sz=3800, center=True):
    algn = "ctr" if center else "l"
    return f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr name="TitleBanner {shape_id}" id="{shape_id}"/>
    <p:cNvSpPr txBox="true"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm rot="0">
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr" rtlCol="false" tIns="0" lIns="0" bIns="0" rIns="0">
      <a:spAutoFit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="{algn}">
        <a:lnSpc><a:spcPts val="4923"/></a:lnSpc>
        <a:spcBef><a:spcPct val="0"/></a:spcBef>
      </a:pPr>
      <a:r>
        <a:rPr lang="en-US" b="true" i="true" sz="{sz}" dirty="0">
          <a:solidFill><a:srgbClr val="{BLUE}"/></a:solidFill>
          <a:latin typeface="{TITLE_FONT}"/>
          <a:ea typeface="{TITLE_FONT}"/>
          <a:cs typeface="{TITLE_FONT}"/>
        </a:rPr>
        <a:t>{_esc(text)}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>"""

# ── Coloured tag/chip ──────────────────────────────────────────────────────────
def _chip_xml(sid, x, y, w, h, text, bg, fg=WHITE, sz=1600, bold=True):
    return f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr name="Chip {sid}" id="{sid}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm rot="0">
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="roundRect">
      <a:avLst><a:gd name="adj" fmla="val 20000"/></a:avLst>
    </a:prstGeom>
    <a:solidFill><a:srgbClr val="{bg}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr" rtlCol="false" tIns="45720" lIns="91440" bIns="45720" rIns="91440"/>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="ctr"/>
      <a:r>
        <a:rPr lang="en-US" b="{'true' if bold else 'false'}" sz="{sz}" dirty="0">
          <a:solidFill><a:srgbClr val="{fg}"/></a:solidFill>
          <a:latin typeface="{BODY_FONT}"/>
        </a:rPr>
        <a:t>{_esc(text)}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>"""

# ── Horizontal divider line ────────────────────────────────────────────────────
def _line_xml(sid, x, y, w, color=BLUE, thickness=18000):
    return f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr name="Line {sid}" id="{sid}"/>
    <p:cNvSpPr><a:spLocks noGrp="true"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm rot="0">
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="0"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="{thickness}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    </a:ln>
  </p:spPr>
</p:sp>"""

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDER HELPERS
# ════════════════════════════════════════════════════════════════════════════════

def new_slide():
    slide = prs.slides.add_slide(BLANK)
    rId = add_bg(slide)
    add_xml(slide, _bg_freeform_xml(rId))
    return slide

def add_section_title(slide, title, subtitle=None):
    """Centered title + optional subtitle — matches ref slides 2,3,5."""
    add_xml(slide, _title_banner_xml(10, 1.0, 0.55, 18.0, 0.80, title,
                                     sz=4100, center=True))
    add_xml(slide, _line_xml(11, 3.0, 1.5, 14.0))
    if subtitle:
        add_xml(slide, _txb_xml(12, 1.0, 1.65, 18.0, 0.65, [{
            "text": subtitle, "sz": 2000, "color": "444444",
            "italic": True, "align": PP_ALIGN.CENTER, "font": BODY_FONT,
        }]))

def add_content_area(slide, x=0.68, y=1.85, w=18.64, h=8.9):
    """White rounded rect — the main content box (matches ref slides 4, 6)."""
    add_xml(slide, _content_rect_xml(20, x, y, w, h))

def bullet_rows(slide, x, y, w, items, sz=1700, color=BLACK, indent=0.35, spacing=0.52):
    """Add bullet point rows."""
    shapes = []
    for i, (bullet, text) in enumerate(items):
        # bullet symbol
        shapes.append(add_xml(slide, _txb_xml(100 + i*2, x, y + i*spacing, indent, spacing, [{
            "text": bullet, "sz": sz, "bold": True, "color": BLUE, "font": BODY_FONT,
        }])))
        # text
        shapes.append(add_xml(slide, _txb_xml(101 + i*2, x+indent, y + i*spacing, w-indent, spacing, [{
            "text": text, "sz": sz, "color": color, "font": BODY_FONT,
        }])))
    return shapes

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()

# Main title banner (matches ref Group 3 style)
add_xml(s, _txb_xml(3, 1.0, 2.7, 18.0, 1.8, [{
    "text": "FlangParallelAnalyzer",
    "sz": 8200, "bold": True, "italic": True,
    "color": BLUE, "font": TITLE_FONT,
    "align": PP_ALIGN.CENTER,
}]))
# Subtitle
add_xml(s, _txb_xml(4, 1.5, 4.65, 17.0, 1.2, [{
    "text": "Automatic Loop Parallelism Analysis for Fortran\nusing LLVM Flang & MLIR",
    "sz": 3600, "bold": True,
    "color": BLACK, "font": BODY_FONT,
    "align": PP_ALIGN.CENTER,
}]))
# Authors
add_xml(s, _txb_xml(5, 1.5, 6.0, 17.0, 0.7, [{
    "text": "Aryan Gupta   ·   Arushi Vaidya   ·   Arpita",
    "sz": 2000, "italic": True, "color": "444444",
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))
add_xml(s, _txb_xml(6, 1.5, 6.65, 17.0, 0.55, [{
    "text": "Dept. of Computer Science and Engineering  |  RV College of Engineering, Bengaluru",
    "sz": 1600, "color": GRAY, "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))
# Stat chips
for i, (num, lbl, bg) in enumerate([
    ("35",    "Test Cases",     "083C92"),
    ("100%",  "Accuracy",       "1A6B3C"),
    ("0",     "False Positives","8B1A1A"),
]):
    cx = 5.0 + i * 3.8
    add_xml(s, _chip_xml(20+i, cx, 7.6, 3.2, 0.65, f"{num}  {lbl}", bg, sz=1800))
# College footer
add_xml(s, _txb_xml(30, 0.5, 10.6, 19.0, 0.45, [{
    "text": "RV COLLEGE OF ENGINEERING",
    "sz": 1600, "color": GRAY, "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Problem Statement")
add_content_area(s)

# Two-column layout inside white box
# Left: context
add_xml(s, _txb_xml(40, 1.1, 2.1, 8.5, 0.55, [{
    "text": "The Challenge", "sz": 2200, "bold": True,
    "color": BLUE, "font": TITLE_FONT,
}]))
bullet_rows(s, 1.1, 2.75, 8.5, [
    ("▸", "Fortran DO loops power scientific computing —\n   weather models, CFD, matrix solvers"),
    ("▸", "OpenMP !$OMP PARALLEL DO enables multi-core speedup"),
    ("▸", "Wrong annotation  →  silent incorrect results"),
    ("▸", "Missed annotation  →  wasted performance"),
    ("▸", "Manual analysis of large codebases = hours of work"),
], sz=1650, color=BLACK, indent=0.3, spacing=0.62)

# Right: today's manual process
add_xml(s, _txb_xml(50, 10.2, 2.1, 8.2, 0.55, [{
    "text": "Manual Process Today", "sz": 2200, "bold": True,
    "color": BLUE, "font": TITLE_FONT,
}]))
steps = [
    "Read the FIR / assembly output",
    "Trace every array subscript for aliases",
    "Check for loop-carried dependencies",
    "Identify reduction variables manually",
    "Write !$OMP directive by hand",
    "Risk of human error at every step",
]
for j, st in enumerate(steps):
    y = 2.75 + j * 0.62
    add_xml(s, _chip_xml(60+j, 10.2, y+0.08, 0.45, 0.38, str(j+1), BLUE, sz=1400))
    add_xml(s, _txb_xml(70+j, 10.75, y, 7.5, 0.55, [{
        "text": st, "sz": 1650, "color": BLACK, "font": BODY_FONT,
    }]))

add_xml(s, _txb_xml(80, 1.1, 10.55, 17.8, 0.48, [{
    "text": "⚠  FlangParallelAnalyzer automates this — per loop, in under a second, with zero false positives",
    "sz": 1700, "bold": True, "italic": True, "color": BLUE,
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Objectives")
add_content_area(s)

add_xml(s, _txb_xml(40, 1.1, 2.0, 17.8, 0.55, [{
    "text": "What FlangParallelAnalyzer sets out to do",
    "sz": 2000, "bold": True, "italic": True,
    "color": BLUE, "font": TITLE_FONT, "align": PP_ALIGN.CENTER,
}]))

obj = [
    ("01", "Automatically classify every Fortran DO loop as SAFE, REDUCTION, or UNSAFE"),
    ("02", "Operate on FIR (Flang Intermediate Representation) — no Fortran source parsing needed"),
    ("03", "Detect loop-carried dependencies using induction-variable index pattern matching"),
    ("04", "Identify scalar reduction variables and emit correct OpenMP REDUCTION clauses"),
    ("05", "Provide a full per-loop analysis trace so developers can audit and override decisions"),
    ("06", "Generate an interactive HTML report with source annotations and confidence scores"),
    ("07", "Guarantee zero false positives — conservatively prefer UNSAFE over incorrect SAFE"),
]

for j, (num, text) in enumerate(obj):
    row = j // 2
    col = j % 2
    x = 1.2 + col * 9.3
    y = 2.8 + row * 1.05
    w = 8.8
    add_xml(s, _chip_xml(50+j, x, y+0.08, 0.6, 0.5, num, BLUE, sz=1400))
    add_xml(s, _txb_xml(60+j, x+0.75, y, w-0.8, 0.85, [{
        "text": text, "sz": 1600, "color": BLACK, "font": BODY_FONT,
    }]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METHODOLOGY (5-PHASE PIPELINE)
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Methodology — 5-Phase Analysis Pipeline",
                  subtitle="Each phase either sets a final verdict or passes to the next")
add_content_area(s)

phases = [
    ("1", "Structure",      "Collect loop bounds,\nnesting depth,\nop count",          "083C92"),
    ("2", "Mem Access",     "Classify reads / writes\nas external or local.\nBuild AccessRecord", "1B6CA8"),
    ("3", "Index Analysis", "Is subscript a(i)\nor a(i±k)?\nDetects loop-carried deps","1A6B3C"),
    ("4", "Reduction",      "Match load → op → store\npattern on scalars.\nEmit REDUCTION clause","B8860B"),
    ("5", "Fallback",       "No ext writes → SAFE\nUnresolved pattern\n→ UNSAFE (conservative)","8B1A1A"),
]

for i, (num, title, body, bg) in enumerate(phases):
    x = 1.1 + i * 3.65
    y_top = 2.05
    # phase number chip
    add_xml(s, _chip_xml(50+i, x+1.2, y_top, 0.6, 0.55, num, bg, sz=2000))
    # phase box
    add_xml(s, _content_rect_xml(60+i, x, y_top+0.65, 3.25, 6.2, fill=LGRAY))
    add_xml(s, _txb_xml(70+i, x+0.1, y_top+0.8, 3.05, 0.6, [{
        "text": title, "sz": 1900, "bold": True,
        "color": bg, "font": TITLE_FONT,
    }]))
    add_xml(s, _line_xml(80+i, x+0.1, y_top+1.52, 3.0, color=bg, thickness=9000))
    add_xml(s, _txb_xml(90+i, x+0.12, y_top+1.7, 3.0, 4.8, [{
        "text": body, "sz": 1550, "color": BLACK, "font": BODY_FONT,
    }]))
    # arrow between boxes
    if i < 4:
        add_xml(s, _txb_xml(200+i, x+3.28, y_top+3.2, 0.37, 0.55, [{
            "text": "→", "sz": 2200, "bold": True, "color": BLUE,
            "align": PP_ALIGN.CENTER, "font": BODY_FONT,
        }]))

# verdict chips at bottom
add_xml(s, _txb_xml(300, 1.1, 9.15, 3.5, 0.45, [{
    "text": "Input: .fir file", "sz": 1500, "italic": True,
    "color": "555555", "font": BODY_FONT,
}]))
for i, (lbl, bg) in enumerate([("SAFE","1A6B3C"),("REDUCTION","B8860B"),("UNSAFE","8B1A1A")]):
    add_xml(s, _chip_xml(310+i, 13.5 + i*2.0, 9.1, 1.8, 0.48, lbl, bg, sz=1500))
add_xml(s, _txb_xml(320, 12.8, 9.1, 1.5, 0.48, [{
    "text": "Output:", "sz": 1500, "bold": True, "color": BLACK, "font": BODY_FONT,
}]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Tech Stack")
add_content_area(s)

tech = [
    ("LLVM 18",        "Compiler infrastructure — pass framework, IR walker API",   "083C92"),
    ("MLIR",           "Multi-Level IR — PassWrapper, op.walk, dyn_cast patterns",  "1B6CA8"),
    ("Flang / FIR",    "Fortran → FIR compiler.  FIR dialect preserves DO loop\nstructure, array shapes, and variable intent", "1A6B3C"),
    ("C++ 17",         "Analysis pass implementation (~700 LOC core logic)",         "444444"),
    ("Python 3",       "Test runner (run_tests.py) and HTML report generator\n(report.py — 1259 LOC)", "B8860B"),
    ("CMake ≥ 3.20",   "Build system.  Builds in < 2 minutes on Codespaces",        "555555"),
]

for i, (name, desc, bg) in enumerate(tech):
    row = i // 2; col = i % 2
    x = 1.1 + col * 9.25
    y = 2.1 + row * 2.5
    add_xml(s, _chip_xml(50+i, x, y, 3.5, 0.55, name, bg, sz=1700))
    add_xml(s, _txb_xml(60+i, x, y+0.68, 8.8, 1.5, [{
        "text": desc, "sz": 1600, "color": BLACK, "font": BODY_FONT,
    }]))

add_xml(s, _txb_xml(80, 1.1, 9.8, 17.8, 0.55, [{
    "text": "Why FIR?  Every array subscript is an explicit SSA value — no need to parse raw Fortran syntax quirks.",
    "sz": 1650, "italic": True, "color": "555555",
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE FLOW
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Architecture Flow")
add_content_area(s)

add_xml(s, _txb_xml(40, 1.1, 2.05, 17.8, 0.5, [{
    "text": "[ Insert Mermaid / Architecture diagram here — see diagram code on next slide ]",
    "sz": 1600, "italic": True, "color": "999999",
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# Component cards
comps = [
    ("LoopParallelAnalysis.cpp\n471 LOC",  "Core MLIR Pass\nPhases 1–5 + printer",      "083C92"),
    ("AccessClassifier.cpp\n165 LOC",      "Phase 2 mem walker\nBase-ref stripping",     "1B6CA8"),
    ("fpa-tool / main.cpp\n71 LOC",        "CLI driver\nDialect registration",           "1A6B3C"),
    ("report.py\n1 259 LOC",               "HTML report generator\nConfidence scoring",  "B8860B"),
]
for i, (name, desc, bg) in enumerate(comps):
    x = 1.1 + i * 4.6
    add_xml(s, _content_rect_xml(50+i, x, 2.75, 4.3, 3.5, fill="EBF2FA"))
    add_xml(s, _chip_xml(60+i, x+0.1, 2.82, 4.1, 0.62, name, bg, sz=1350))
    add_xml(s, _txb_xml(70+i, x+0.15, 3.6, 4.0, 2.3, [{
        "text": desc, "sz": 1550, "color": BLACK, "font": BODY_FONT,
    }]))

# Pipeline flow labels
flow = [".f90 source","FIR (.fir)","MLIR Module","Per-loop Verdict"]
flow_colors = ["083C92","1B6CA8","1A6B3C","B8860B"]
for i, (lbl, bg) in enumerate(zip(flow, flow_colors)):
    x = 1.1 + i * 4.6
    add_xml(s, _chip_xml(80+i, x, 6.55, 4.3, 0.5, lbl, bg, sz=1450))
    if i < 3:
        add_xml(s, _txb_xml(90+i, x+4.3, 6.65, 0.3, 0.4, [{
            "text": "→", "sz": 1800, "bold": True, "color": BLUE,
            "align": PP_ALIGN.CENTER, "font": BODY_FONT,
        }]))

add_xml(s, _txb_xml(100, 1.1, 7.35, 17.8, 0.45, [{
    "text": "flang-new -fc1 -emit-fir  →  fpa-tool input.fir  →  Stdout + HTML report",
    "sz": 1600, "italic": True, "color": "555555",
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TEST RESULTS
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Test Results")
add_content_area(s)

# Summary row
for i, (lbl, n, bg) in enumerate([
    ("SAFE",             "13", "1A6B3C"),
    ("REDUCTION",        " 5", "B8860B"),
    ("UNSAFE (dep)",     "10", "8B1A1A"),
    ("UNSAFE (conserv)", " 7", "6B3A8B"),
]):
    x = 1.1 + i * 4.55
    add_xml(s, _chip_xml(50+i, x, 2.05, 4.3, 0.68, f"{n}   {lbl}", bg, sz=1700))

# Table headers
headers = ["Fortran Pattern", "Expected", "Got", "✓/✗"]
col_x = [1.1, 11.5, 14.6, 17.7]
col_w = [10.1, 2.8, 2.8, 1.8]
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_xml(s, _chip_xml(60+j, cx, 2.95, cw, 0.48, hdr, "083C92", sz=1550))

rows = [
    ("b(i) = a(i) * 2.0",           "SAFE",      "SAFE",      "✓", "1A6B3C"),
    ("total += a(i)*b(i)",           "REDUCTION", "REDUCTION", "✓", "1A6B3C"),
    ("a(i) = a(i) + a(i-1)",         "UNSAFE",    "UNSAFE",    "✓", "1A6B3C"),
    ("b(idx(i)) = a(i)  [scatter]",  "UNSAFE",    "UNSAFE",    "✓", "1A6B3C"),
    ("a(i) = a(i) * 3.0  [inplace]", "UNSAFE",    "UNSAFE",    "✓", "1A6B3C"),
    ("norm2 += a(i)*a(i)",           "REDUCTION", "REDUCTION", "✓", "1A6B3C"),
    ("read-only traversal",          "SAFE",      "SAFE",      "✓", "1A6B3C"),
]
for i, (pat, exp, got, tick, tc) in enumerate(rows):
    y = 3.56 + i * 0.72
    bg_r = "F0F0F0" if i%2==0 else WHITE
    cells = [(pat,"444444"), (exp,"444444"), (got,"444444"), (tick, tc)]
    for j, ((txt, tc2), cx, cw) in enumerate(zip(cells, col_x, col_w)):
        add_xml(s, _content_rect_xml(80+i*4+j, cx, y, cw, 0.62, fill=bg_r))
        bold = j == 3
        add_xml(s, _txb_xml(200+i*4+j, cx+0.08, y+0.08, cw-0.1, 0.5, [{
            "text": txt, "sz": 1500, "bold": bold, "color": tc2, "font": BODY_FONT,
        }]))

add_xml(s, _txb_xml(300, 1.1, 8.65, 17.8, 0.48, [{
    "text": "30 / 30 comprehensive tests pass  ·  5 / 5 original tests pass  ·  35 total  ·  0 false positives",
    "sz": 1600, "bold": True, "italic": True, "color": BLUE,
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LIMITATIONS & FUTURE WORK
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_section_title(s, "Limitations & Future Work")
add_content_area(s)

# Left — limitations
add_xml(s, _txb_xml(40, 1.1, 2.1, 9.0, 0.55, [{
    "text": "Current Limitations",
    "sz": 2200, "bold": True, "italic": True, "color": BLUE, "font": TITLE_FONT,
}]))
lims = [
    "a(i) = a(i)*k  →  conservatively UNSAFE\n(same-index in-place update)",
    "Multi-dim subscripts c(i,j) not traced\n(outer loop conservatively UNSAFE)",
    "Function call side-effects not analyzed",
    "min / max reductions not yet detected\n(only + and × supported)",
    "Chained reductions not matched:\ns = s + f(x) + g(y)",
]
for j, txt in enumerate(lims):
    add_xml(s, _txb_xml(50+j, 1.1, 2.78+j*1.3, 0.35, 0.55, [{
        "text": "–", "sz": 1700, "bold": True, "color": "8B1A1A", "font": BODY_FONT,
    }]))
    add_xml(s, _txb_xml(60+j, 1.55, 2.78+j*1.3, 8.5, 1.15, [{
        "text": txt, "sz": 1600, "color": BLACK, "font": BODY_FONT,
    }]))

# Right — future work
add_xml(s, _txb_xml(70, 10.5, 2.1, 8.8, 0.55, [{
    "text": "Future Work",
    "sz": 2200, "bold": True, "italic": True, "color": BLUE, "font": TITLE_FONT,
}]))
future = [
    "Read-before-write check for in-place updates",
    "Multi-dim subscript tracing per axis",
    "Add min / max reduction patterns",
    "Interprocedural side-effect summaries",
    "Auto-insert !$OMP directives in source",
    "CI / CD pre-commit hook integration",
]
for j, txt in enumerate(future):
    add_xml(s, _txb_xml(80+j, 10.5, 2.78+j*1.12, 0.4, 0.55, [{
        "text": "→", "sz": 1700, "bold": True, "color": BLUE, "font": BODY_FONT,
    }]))
    add_xml(s, _txb_xml(90+j, 11.0, 2.78+j*1.12, 8.2, 0.95, [{
        "text": txt, "sz": 1600, "color": BLACK, "font": BODY_FONT,
    }]))

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — THANK YOU  (matches ref slide 9 exactly)
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()

add_xml(s, _txb_xml(3, 1.5, 3.5, 17.0, 2.5, [{
    "text": "Thank you",
    "sz": 7200, "bold": True, "italic": True,
    "color": BLUE, "font": TITLE_FONT,
    "align": PP_ALIGN.CENTER,
}]))
add_xml(s, _txb_xml(4, 3.0, 6.2, 14.0, 0.7, [{
    "text": "Aryan Gupta   ·   Arushi Vaidya   ·   Arpita",
    "sz": 2000, "color": "444444", "font": BODY_FONT,
    "align": PP_ALIGN.CENTER,
}]))
add_xml(s, _line_xml(5, 4.0, 7.0, 12.0, color=BLUE, thickness=12000))
add_xml(s, _txb_xml(6, 1.5, 7.25, 17.0, 0.5, [{
    "text": "github.com/AryanGupta21/FlangParallelAnalyzer",
    "sz": 1700, "italic": True, "color": GRAY,
    "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))
# College tag at bottom (matches ref)
add_xml(s, _txb_xml(10, 5.5, 10.55, 9.3, 0.48, [{
    "text": "RV COLLEGE OF ENGINEERING",
    "sz": 1600, "color": GRAY, "font": BODY_FONT, "align": PP_ALIGN.CENTER,
}]))

# ════════════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}  ({len(prs.slides)} slides)")
